import os
import random
import numpy as np
import torch
import torch.backends.cudnn as cudnn
import neurokit2 as nk
import matplotlib.pyplot as plt
import biosppy
import scipy
import scipy.signal
import scipy.ndimage
from pantom import Pan_Tompkins_Plus_Plus
from pptx import Presentation
from pptx.util import Inches
import glob
import os
from tqdm import tqdm
from sklearn.metrics import accuracy_score, precision_score, recall_score, f1_score, matthews_corrcoef



def fix_seed(SEED=42):
    os.environ['SEED'] = str(SEED)
    os.environ['CUDA_DEVICE_ORDER'] = 'PCI_BUS_ID'
    np.random.seed(SEED)
    random.seed(SEED)
    torch.manual_seed(SEED)
    torch.cuda.manual_seed(SEED)
    torch.cuda.manual_seed_all(SEED)
    cudnn.benchmark = False
    cudnn.deterministic = True
    random.seed(SEED)
    
def pantom_rpeak_detection(ecg, low_cutoff=5, high_cutoff=15, fs=500):
    """
    ECG 신호에서 R-peak를 검출하는 함수
    
    Args:
        ecg (numpy.ndarray): ECG 신호 데이터
        fs (int): 샘플링 주파수 (Hz)
        
    Returns:
        numpy.ndarray: 검출된 R-peak의 인덱스
    """
    # 파라미터 설정
    max_QRS_duration = 0.150  # sec
    window_size = int(max_QRS_duration * fs)
    
    # 밴드패스 필터 적용
    lowpass = scipy.signal.butter(1, high_cutoff / (fs / 2.0), "low")
    highpass = scipy.signal.butter(1, low_cutoff / (fs / 2.0), "high")
    ecg_low = scipy.signal.filtfilt(*lowpass, x=ecg)
    ecg_band = scipy.signal.filtfilt(*highpass, x=ecg_low)
    
    # 미분 및 제곱
    diff = np.diff(ecg_band)
    squared = np.square(diff)
    
    # 이동 평균 필터
    mwa = np.pad(squared, (window_size - 1, 0), "constant", constant_values=(0, 0))
    mwa = np.convolve(mwa, np.ones(window_size), "valid")
    for i in range(1, window_size):
        mwa[i - 1] = mwa[i - 1] / i
    mwa[window_size - 1 :] = mwa[window_size - 1 :] / window_size
    mwa[: int(max_QRS_duration * fs * 2)] = 0
    
    # 가우시안 필터로 신호 스무딩 및 미분
    energy = scipy.ndimage.gaussian_filter1d(mwa, fs / 8.0)
    energy_diff = np.diff(energy)
    
    # 미분값이 0을 교차하는 지점 찾기 (피크 위치)
    r_peaks = (energy_diff[:-1] > 0) & (energy_diff[1:] < 0)
    r_peaks = np.flatnonzero(r_peaks)
    r_peaks -= int(window_size / 2)
    return {'rpeaks': r_peaks}

def normalize_to_range(signal, min_val, max_val):
    """
    신호를 지정된 범위로 정규화합니다.
    """
    signal_min = np.min(signal)
    signal_max = np.max(signal)
    
    # 신호의 범위가 0인 경우(모든 값이 동일한 경우) 처리
    if signal_max == signal_min:
        return np.zeros_like(signal)
    
    # 신호를 0-1 범위로 정규화한 후 원하는 범위로 스케일링
    normalized = (signal - signal_min) / (signal_max - signal_min)
    scaled = normalized * (max_val - min_val) + min_val
    
    return scaled

def preprocess_ecg(ecg_signal, sampling_rate, lowcut=0.5, highcut=40, min_val=-1, max_val=1, power_line_freq=60):
    if len(ecg_signal) == 0:
        return np.array([])
    
    if np.all(np.isnan(ecg_signal)):
        return np.zeros_like(ecg_signal)
    
    pad_length = int(sampling_rate)
    padded_signal = np.pad(ecg_signal, (pad_length, pad_length), 'symmetric')
    
    baseline = scipy.signal.savgol_filter(padded_signal, int(sampling_rate/2)*2+1, 3)
    baseline_removed = padded_signal - baseline
    
    # 버터워스 필터 적용
    nyquist = 0.5 * sampling_rate
    order = 2
    low = lowcut / nyquist
    high = highcut / nyquist
    b, a = scipy.signal.butter(order, [low, high], btype='band')
    filtered_signal = scipy.signal.filtfilt(b, a, baseline_removed)
    
    # 노치 필터는 선택적으로 적용
    if power_line_freq > 0:
        w0 = power_line_freq / nyquist
        Q = 30.0
        b_notch, a_notch = scipy.signal.iirnotch(w0, Q, sampling_rate)
        filtered_signal = scipy.signal.filtfilt(b_notch, a_notch, filtered_signal)
    
    # 패딩 제거
    filtered_signal = filtered_signal[pad_length:-pad_length]
    
    # 정규화
    min_signal = np.min(filtered_signal)
    max_signal = np.max(filtered_signal)
    if max_signal > min_signal:
        processed_signal = (filtered_signal - min_signal) / (max_signal - min_signal)
        processed_signal = processed_signal * (max_val - min_val) + min_val
    else:
        processed_signal = np.zeros_like(filtered_signal)
    
    return processed_signal

def ensemble_rpeaks_detection(all_rpeaks, preprocessed_signal, sampling_rate, min_votes=3, use_weighted=True, debug=False, algorithm_weights=None, algorithm_names=None):
    """
    여러 알고리즘에서 검출된 R-peaks를 앙상블하여 최종 R-peaks를 결정하는 함수
    
    Parameters:
    -----------
    all_rpeaks : list
        각 알고리즘에서 검출된 R-peaks 리스트
    preprocessed_signal : ndarray
        전처리된 ECG 신호
    sampling_rate : int
        샘플링 레이트 (Hz)
    min_votes : int, optional
        최소 투표 수 (기본값: 3)
    use_weighted : bool, optional
        가중치 기반 앙상블 사용 여부 (기본값: True)
    debug : bool, optional
        디버깅 정보 출력 여부 (기본값: False)
    algorithm_weights : dict, optional
        각 알고리즘별 사전 정의된 가중치 (기본값: None)
    algorithm_names : list, optional
        all_rpeaks 리스트와 매칭되는 알고리즘 이름 리스트 (기본값: None)
        
    Returns:
    --------
    ndarray
        앙상블 방법으로 보정된 최종 R-peaks
    """
    # 모든 R-peaks를 하나의 배열로 합치고 정렬
    all_peaks_combined = np.sort(np.concatenate(all_rpeaks))
    
    # 유사한 위치의 피크를 그룹화하기 위한 허용 오차 설정 (샘플 단위)
    tolerance = int(0.05 * sampling_rate)  # 50ms 허용 오차

    # 그룹화된 피크 위치 저장 (최적화된 방식)
    grouped_peaks = []
    current_group = [all_peaks_combined[0]]
    
    for peak in all_peaks_combined[1:]:
        if peak - current_group[-1] <= tolerance:
            current_group.append(peak)
        else:
            grouped_peaks.append(int(np.median(current_group)))
            current_group = [peak]
            
    # 마지막 그룹 처리
    if current_group:
        grouped_peaks.append(int(np.median(current_group)))
    
    # 가중치 계산 (use_weighted가 False면 건너뜀)
    weights = np.ones(len(all_rpeaks))
    
    if use_weighted:
        # 1. RR 간격 기반 품질 점수 계산
        for i, algorithm_peaks in enumerate(all_rpeaks):
            if len(algorithm_peaks) > 1:
                # RR 간격 계산 (ms 단위)
                rr_intervals = np.diff(algorithm_peaks) / sampling_rate * 1000
                
                # 생리학적 유효성 검사를 벡터화
                valid_mask = (rr_intervals >= 300) & (rr_intervals <= 2000)
                valid_ratio = np.mean(valid_mask)
                
                # 통계 계산
                rr_mean = np.mean(rr_intervals)
                rr_std = np.std(rr_intervals)
                rr_median = np.median(rr_intervals)
                
                # 변동 계수 (낮을수록 좋음)
                cv = rr_std / rr_mean if rr_mean > 0 else 1.0
                
                # 극단적 이상치 비율 (낮을수록 좋음)
                outlier_mask = np.abs(rr_intervals - rr_median) > (3 * rr_std)
                outlier_ratio = np.mean(outlier_mask)
                
                # 심박수 범위 검사 (60-180 BPM 범위 내 비율)
                hr_valid_mask = (rr_intervals >= 333) & (rr_intervals <= 1000)  # 60-180 BPM
                hr_valid_ratio = np.mean(hr_valid_mask)
                
                # 종합 품질 점수 계산 (0~1 사이 값) - 심박수 범위 검사 추가
                quality_score = (valid_ratio * 0.4) + ((1 - cv) * 0.2) + ((1 - outlier_ratio) * 0.2) + (hr_valid_ratio * 0.2)
                weights[i] = max(0.1, min(1.0, quality_score))
            else:
                weights[i] = 0.1
        
        # 2. 사전 정의된 알고리즘 가중치 적용 (있는 경우)
        if algorithm_weights is not None and algorithm_names is not None:
            for i, alg_name in enumerate(algorithm_names):
                if i < len(weights) and alg_name in algorithm_weights:
                    # 품질 점수와 사전 정의된 가중치를 곱하여 최종 가중치 계산
                    weights[i] *= algorithm_weights[alg_name]
        
        # 가중치 정규화 (합이 알고리즘 수와 같도록)
        weights = weights / np.sum(weights) * len(weights)
        
        # 디버깅 정보 출력
        if debug and algorithm_names is not None:
            print("알고리즘별 가중치:")
            for i, alg_name in enumerate(algorithm_names):
                if i < len(weights):
                    print(f"  {alg_name}: {weights[i]:.2f}")
    
    # 투표 계산 (딕셔너리 대신 defaultdict 사용)
    from collections import defaultdict
    votes = defaultdict(float)
    
    # 각 알고리즘의 R-peaks에 대해 투표 계산
    for i, algorithm_peaks in enumerate(all_rpeaks):
        weight = weights[i]
        
        # 각 피크에 대해 가장 가까운 그룹화된 피크 찾기
        for peak in algorithm_peaks:
            # NumPy 배열 연산으로 최적화
            distances = np.abs(np.array(grouped_peaks) - peak)
            closest_idx = np.argmin(distances)
            closest_peak = grouped_peaks[closest_idx]
            
            # 허용 오차 내에 있는지 확인
            if distances[closest_idx] <= tolerance:
                votes[closest_peak] += weight
                
    # 최소 투표 수 이상을 받은 피크만 선택
    ensemble_rpeaks = np.array([peak for peak, vote_count in votes.items() if vote_count >= min_votes])
    ensemble_rpeaks = np.sort(ensemble_rpeaks)
    
    # 디버깅 정보 출력
    if debug:
        from pprint import pprint
        sorted_votes = sorted(votes.items(), key=lambda x: x[1], reverse=True)
        print(f"최소 투표 수: {min_votes:.2f}")
        print("투표 수가 높은 순서대로 정렬:")
        pprint(sorted_votes)
    
    return ensemble_rpeaks

def correct_rpeaks(signal=None, rpeaks=None, sampling_rate=1000.0, tol=0.05):
    """Correct R-peak locations to the maximum within a tolerance.

    Parameters
    ----------
    signal : array
        ECG signal.
    rpeaks : array
        R-peak location indices.
    sampling_rate : int, float, optional
        Sampling frequency (Hz).
    tol : int, float, optional
        Correction tolerance (seconds).

    Returns
    -------
    rpeaks : array
        Cerrected R-peak location indices.

    Notes
    -----
    * The tolerance is defined as the time interval :math:`[R-tol, R+tol[`.

    """
    # 입력 확인
    if signal is None:
        raise TypeError("Please specify an input signal.")

    if rpeaks is None:
        raise TypeError("Please specify the input R-peaks.")

    tol = int(tol * sampling_rate)
    length = len(signal)

    newR = []
    for r in rpeaks:
        a = r - tol
        if a < 0:
            a = 0
        b = r + tol
        if b > length:
            b = length
        newR.append(a + np.argmax(signal[a:b]))

    newR = sorted(list(set(newR)))
    newR = np.array(newR, dtype="int")

    return newR


def plot_ecg_with_multiple_algorithms(results, title, lead='I', plot_path=None, signal_quality=None):
    """
    여러 알고리즘의 R-peak 검출 결과를 하나의 그래프로 비교하는 함수
    
    Parameters:
    -----------
    results : dict
        extract_ecg_and_rpeaks 함수에서 반환된 결과 딕셔너리
    title : str
        그래프 제목
    lead : str
        ECG 리드 이름
    plot_path : str, optional
        그래프를 저장할 경로
    signal_quality : float, optional
        신호 품질 점수
    """
    # 필요한 데이터 추출
    ensemble_rpeaks = results.get('ensemble_rpeaks', [])
    algorithm_results = results.get('algorithm_results', {})
    emsemble_plot_signal = results.get('plot_signal', [])
    
    # 서브플롯 설정
    fig, axes = plt.subplots(len(algorithm_results) + 1, 1, figsize=(20, 20))
    fig.suptitle(f"{title} \n(Signal Quality: {signal_quality:.2f})", fontsize=16)
    
    # 각 알고리즘별 결과 플롯
    for i, (name, result) in enumerate(algorithm_results.items()):
        ax = axes[i]
        plot_signal = result.get('plot_signal', [])
        ax.plot(plot_signal, label="Signal")

        # R-peaks 표시
        r_peak = result.get('original_rpeaks', [])
        corrected_rpeaks = result.get('corrected_rpeaks', [])
        
        if len(r_peak) > 0:
            ax.scatter(r_peak, plot_signal[r_peak], 
                      color='blue', marker='o', s=50, label='R-peaks')        
        if len(corrected_rpeaks) > 0:
            ax.scatter(corrected_rpeaks, plot_signal[corrected_rpeaks], 
                      color='red', marker='o', s=50, label='Corr. R-peaks')
                
        ax.set_title(f'{name}')
        ax.set_ylabel('Amplitude')
        ax.grid(True, alpha=0.3)
        ax.legend(loc='upper left')
    
    # 앙상블 결과 플롯
    ax_ensemble = axes[-1]
    ax_ensemble.plot(emsemble_plot_signal, label="Signal")
    
    if len(ensemble_rpeaks) > 0:
        ax_ensemble.scatter(ensemble_rpeaks, emsemble_plot_signal[ensemble_rpeaks], 
                           color='purple', marker='*', s=50, label='Ensemble')
            
    ax_ensemble.set_title('Ensemble Method (R-peaks Detection)')
    ax_ensemble.set_ylabel('Amplitude')
    ax_ensemble.set_xlabel('Samples')
    ax_ensemble.grid(True, alpha=0.3)
    ax_ensemble.legend(loc='upper left')
    
    plt.tight_layout()
    plt.subplots_adjust(top=0.9)  # 제목 공간 확보
    
    if plot_path:
        plt.savefig(plot_path)
    else:
        plt.show()
    plt.close()
    
def extract_ecg_and_rpeaks(raw_signal, sampling_rate, lowcut=0.5, highcut=40, min_votes_ratio=0.3, debug=False, r_min=250, r_max=4750):
    """
    ECG 신호에서 R-peak를 추출하는 함수
    앙상블 방법을 사용하여 더 정확한 R-peak 검출
    
    Parameters:
    -----------
    raw_signal : numpy.ndarray
        ECG 신호 데이터
    sampling_rate : float
        샘플링 레이트
    lowcut : float
        대역 통과 필터의 하한 주파수 (Hz)
    highcut : float
        대역 통과 필터의 상한 주파수 (Hz)
    min_votes_ratio : float
        최소 투표 비율 (0~1 사이, 기본값: 0.3)
        전체 알고리즘 중 이 비율 이상의 알고리즘이 동의해야 R-peak로 인정
    debug : bool
        debug 모드 여부
    r_min : int, optional
        R-peak 검출 시작 인덱스 (기본값: 250)
    r_max : int, optional
        R-peak 검출 종료 인덱스 (기본값: 4750)
        
    Returns:
    --------
    dict: {
        'raw_signal': 원본 ECG 신호,
        'preprocessed_signal': 통일된 전처리 ECG 신호,
        'ensemble_rpeaks': 앙상블 방법으로 결정된 R-peak 위치,
        'plot_signal': 시각화를 위한 신호,
        'algorithm_results': 각 알고리즘별 R-peak 결과 딕셔너리,
        'signal_quality': 신호 품질 점수
    }
    """
    m_detector = Pan_Tompkins_Plus_Plus()
    
    # 신호 품질 평가
    signal_quality = assess_signal_quality(raw_signal, sampling_rate)
    
    # 전처리 신호 생성 - 모든 알고리즘에 동일하게 사용
    preprocessed_signal = preprocess_ecg(raw_signal, sampling_rate, lowcut, highcut)
    
    info_pantom = pantom_rpeak_detection(raw_signal, sampling_rate)
    info_pantom_pp = m_detector.rpeak_detection(raw_signal, sampling_rate)
    
    ecg_cleaned = nk.ecg_clean(raw_signal, sampling_rate=sampling_rate)
    _, info_nk = nk.ecg_peaks(ecg_cleaned=ecg_cleaned, sampling_rate=sampling_rate)
    
    algorithms = {
        'NeuroKit2': {'rpeaks': info_nk['ECG_R_Peaks']},
        'Pan-Tomkins': {'rpeaks': info_pantom['rpeaks']},
        'Pan-Tomkins++': {'rpeaks': info_pantom_pp['rpeaks']},
        'Hamilton': biosppy.signals.ecg.hamilton_segmenter,
        'Christov': biosppy.signals.ecg.christov_segmenter,
        'Engzee': biosppy.signals.ecg.engzee_segmenter
    }
    
    num_algorithms = len(algorithms)
    min_votes = max(2, round(num_algorithms * min_votes_ratio))
    
    # 각 알고리즘 적용 및 보정
    all_rpeaks = []
    algorithm_results = {}
    
    for name, algorithm in algorithms.items():
        if isinstance(algorithm, dict):
            rpeaks = algorithm['rpeaks']
        else:
            result = algorithm(preprocessed_signal, sampling_rate=sampling_rate)
            rpeaks = result['rpeaks']
            
        # 모든 알고리즘에 동일한 신호로 R-peak 위치 보정
        corrected_rpeaks = correct_rpeaks(raw_signal, rpeaks, sampling_rate=sampling_rate, tol=0.05)
        all_rpeaks.append(corrected_rpeaks)
        
        # 각 알고리즘별 결과 저장
        algorithm_results[name] = {
            'original_rpeaks': rpeaks,
            'corrected_rpeaks': corrected_rpeaks,
            'plot_signal': raw_signal  # 모든 알고리즘에 동일한 신호 사용
        }
    
    # 알고리즘 가중치 설정
    algorithm_weights = {
        'NeuroKit2': 1.0,
        'Pan-Tomkins': 0.8,
        'Pan-Tomkins++': 1.3,
        'Hamilton': 0.8,
        'Christov': 1.2,
        'Engzee': 0.7,
    }
    
    # 신호 품질에 따른 가중치 조정
    if signal_quality < 0.7:  # 낮은 품질
        if debug:
            print(f"신호 품질이 매우 낮음 ({signal_quality:.2f})")
        # 노이즈에 강한 알고리즘 우선
        for alg in algorithm_weights:
            if alg == 'Engzee':
                algorithm_weights[alg] *= 2.0
            else:
                algorithm_weights[alg] *= 0.9
    elif signal_quality < 0.75:  # 보통 품질
        if debug:
            print(f"신호 품질이 낮음 ({signal_quality:.2f})")
        algorithm_weights['Engzee'] *= 1.5
        algorithm_weights['Christov'] *= 0.9
    elif signal_quality > 0.85:  # 높은 품질
        if debug:
            print(f"신호 품질이 우수함 ({signal_quality:.2f})")
        # 정밀한 알고리즘 우선
        algorithm_weights['NeuroKit2'] *= 1.2
        algorithm_weights['Hamilton'] *= 1.1
        algorithm_weights['Pan-Tomkins++'] *= 1.1
        
    # 알고리즘 이름 리스트 (all_rpeaks와 동일한 순서)
    algorithm_names = list(algorithm_results.keys())
    
    # 앙상블 방법으로 최종 R-peaks 결정 (가중치 전달)
    ensemble_rpeaks = ensemble_rpeaks_detection(
        all_rpeaks, 
        raw_signal,
        sampling_rate, 
        min_votes=min_votes, 
        use_weighted=True,
        algorithm_weights=algorithm_weights,
        algorithm_names=algorithm_names,
        debug=debug 
    )
    ensemble_rpeaks = ensemble_rpeaks[(ensemble_rpeaks >= r_min) & (ensemble_rpeaks <= r_max)]
    # 결과를 딕셔너리로 반환 - 키 이름 통일
    results = {
        'raw_signal': raw_signal,
        'preprocessed_signal': preprocessed_signal,
        'ensemble_rpeaks': ensemble_rpeaks,
        'plot_signal': raw_signal,
        'algorithm_results': algorithm_results,
        'signal_quality': signal_quality
    }
    return results


def plot_ecg_with_peaks(signal, signals_info, title, show_r_peaks=True, show_other_peaks=False, show_wave_markers=False, path=None):
    """
    ECG 신호와 피크를 시각화하는 함수
    
    Parameters:
    -----------
    signal : numpy.ndarray
        ECG 신호 데이터
    signals_info : dict
        ECG 신호의 피크 정보를 담은 딕셔너리
    patient_index : int
        환자 인덱스
    show_r_peaks : bool
        R-peaks 표시 여부
    show_other_peaks : bool
        R-peaks 외 다른 피크 표시 여부
    show_wave_markers : bool
        파형 마커(onsets, offsets) 표시 여부
    path : str, optional
        그래프를 저장할 경로
    """
    plt.figure(figsize=(20, 10))
    
    # 정제된 ECG 신호 플롯
    plt.plot(signal, label="Raw ECG")
    # 피크 타입 정의
    peak_types = {
        'ECG_R_Peaks': {'color': 'red', 'marker': 'o', 'label': 'R-peaks'},
    }
    
    # 다른 피크 타입 추가 (show_other_peaks가 True일 때만)
    if show_other_peaks:
        peak_types.update({
            'ECG_T_Peaks': {'color': 'purple', 'marker': 'o', 'label': 'T-peaks'},
            'ECG_P_Peaks': {'color': 'orange', 'marker': 'o', 'label': 'P-peaks'},
            'ECG_Q_Peaks': {'color': 'green', 'marker': 'o', 'label': 'Q-peaks'},
            'ECG_S_Peaks': {'color': 'blue', 'marker': 'o', 'label': 'S-peaks'}
        })
    
    # R-peaks만 표시하거나 show_r_peaks가 True일 때
    if show_r_peaks:
        # 모든 피크 타입 플롯 - 정제된 신호에 맞춰 표시
        for peak_name, peak_info in peak_types.items():
            if peak_name in signals_info.keys():
                peak_indices = signals_info[peak_name]
                if len(peak_indices) > 0:
                    plt.scatter(
                        peak_indices, 
                        signal[peak_indices],
                        color=peak_info['color'], 
                        marker=peak_info['marker'], 
                        label=peak_info['label']
                    )

    # 시작점과 종료점 마커 정의 (show_wave_markers가 True일 때만)
    if show_wave_markers:
        print("signals_info : ", signals_info)
        print(signals_info['ECG_P_Onsets'].shape)
        
        wave_markers = {
            # Onsets
            'ECG_P_Onsets': {'color': 'blue', 'linestyle': '--', 'label': 'P-onsets'},
            'ECG_T_Onsets': {'color': 'cyan', 'linestyle': '--', 'label': 'T-onsets'},
            'ECG_QRS_Onsets': {'color': 'magenta', 'linestyle': '--', 'label': 'QRS-onsets'},
            # Offsets
            'ECG_P_Offsets': {'color': 'blue', 'linestyle': ':', 'label': 'P-offsets'},
            'ECG_T_Offsets': {'color': 'cyan', 'linestyle': ':', 'label': 'T-offsets'},
            'ECG_QRS_Offsets': {'color': 'magenta', 'linestyle': ':', 'label': 'QRS-offsets'}
        }

        for marker_name, marker_info in wave_markers.items():
            if marker_name in signals_info.keys():
                indices = np.where(signals_info[marker_name] == 1)[0]
                if len(indices) > 0:
                    plt.axvline(x=indices[0], color=marker_info['color'], 
                                linestyle=marker_info['linestyle'], alpha=0.7, 
                                label=marker_info['label'])
                    
                    # 나머지는 범례 없이 표시
                    for idx in indices[1:]:
                        plt.axvline(x=idx, color=marker_info['color'], 
                                    linestyle=marker_info['linestyle'], alpha=0.7)
    plt.title(f'ECG Signal Analysis with Detected Peaks')
    plt.xlabel('Samples')
    plt.ylabel('Amplitude')
    plt.grid(True, alpha=0.3)
    plt.legend(loc='upper right', bbox_to_anchor=(1.15, 1))  # 범례 위치 조정
    plt.tight_layout()  # 레이아웃 자동 조정
    if path:
        plt.savefig(path)
    else:
        plt.show()
        
    plt.close()
    
def find_peaks_relative_to_rr(signal, r_peaks, sampling_rate=500, debug=False, tol=0.2):
    """
    절대 시간 범위(ms)와 상대 범위(%)를 조합하여 R-peak 기준으로 다른 피크들을 찾는 함수
    
    Args:
        signal: ECG 신호
        r_peaks: 검출된 R-peak 위치 배열
        sampling_rate: 샘플링 레이트 (Hz)
        debug: 디버그 모드 활성화 여부
        tol: 피크 검출 시 허용 오차 (0.0~1.0 사이의 값, 기본값: 0.0)
            - 값이 클수록 더 넓은 범위에서 피크를 검색함
        
    Returns:
        peaks: 모든 피크 위치를 담은 딕셔너리
    """
    peaks = {'ECG_R_Peaks': r_peaks, 'ECG_Q_Peaks': [], 'ECG_S_Peaks': [], 'ECG_P_Peaks': [], 'ECG_T_Peaks': []}
    
    # 임상적 절대 시간 범위 (ms)
    clinical_p_range = (-240, -80)
    clinical_q_range = (-40, -30)
    clinical_s_range = (20, 50)
    clinical_t_range = (150, 300)
    
    # RR 간격을 기준으로 탐색 범위 동적 조정
    if len(r_peaks) > 1:
        rr_intervals = np.diff(r_peaks) / sampling_rate * 1000
        mean_rr = np.mean(rr_intervals)
    else:
        # R-peak가 하나만 있는 경우 기본값 설정
        mean_rr = 1000  # 1초 (60 bpm에 해당)
    
    # 심박수 계산 (bpm)
    hr = 60000 / mean_rr
    if debug:
        print(f"평균 RR 간격: {mean_rr:.2f}ms, 심박수: {hr:.2f}bpm")
    
    # 정상 심박수 범위 정의
    normal_hr_min = 60
    normal_hr_max = 100
    
    # 샘플 변환 계수 미리 계산
    ms_to_sample = sampling_rate / 1000
    
    # 각 R-peak에 대해 개별적으로 탐색 범위 조정
    for i, r_idx in enumerate(r_peaks):
        # 현재 심박의 RR 간격 계산
        if i > 0:
            # 이전 RR 간격 (현재 R-peak와 이전 R-peak 사이)
            pre_local_rr = (r_idx - r_peaks[i-1]) / ms_to_sample  # ms
        else:
            # 첫 번째 R-peak는 다음 간격 사용 (없으면 평균 사용)
            pre_local_rr = mean_rr
            
        # 다음 RR 간격 (현재 R-peak와 다음 R-peak 사이)
        if i < len(r_peaks) - 1:
            post_local_rr = (r_peaks[i+1] - r_idx) / ms_to_sample  # ms
            next_r_idx = r_peaks[i+1]
            estimated_q_offset = round(clinical_q_range[0] * ms_to_sample)  # Q-peak의 일반적인 오프셋
            next_q_idx = next_r_idx + estimated_q_offset  # 다음 R-peak 기준 Q-peak 예상 위치
        else:
            # 마지막 R-peak는 이전 간격 사용
            post_local_rr = pre_local_rr
            next_r_idx = float('inf')  # 기본값으로 무한대 설정
            next_q_idx = float('inf')  # 기본값으로 무한대 설정
        
        # 심박수 계산 (이전 및 다음 RR 간격 기반)
        pre_local_hr = 60000 / pre_local_rr
        post_local_hr = 60000 / post_local_rr
        
        # 심박수에 따른 동적 조정 계수 계산
        pre_local_hr_factor = 1.0
        post_local_hr_factor = 1.0
        
        if pre_local_hr <= normal_hr_min:
            # 느린 심박수: 정상 최소값(60bpm)을 기준으로 조정
            pre_local_hr_factor = normal_hr_min / pre_local_hr
        elif pre_local_hr >= normal_hr_max:
            # 빠른 심박수: 정상 최대값(100bpm)을 기준으로 조정
            pre_local_hr_factor = normal_hr_max / pre_local_hr
            
        if post_local_hr <= normal_hr_min:
            post_local_hr_factor = normal_hr_min / post_local_hr
        elif post_local_hr >= normal_hr_max:
            post_local_hr_factor = normal_hr_max / post_local_hr
        
        # 절대 범위 조정 (심박수에 따라 약간 조정)
        # P, Q 파형은 이전 심박수 기준으로 조정
        # S, T 파형은 다음 심박수 기준으로 조정
        
        # 샘플 단위로 변환 (한 번에 계산)
        p_samples = (round(clinical_p_range[0] * pre_local_hr_factor * ms_to_sample), 
                     round(clinical_p_range[1] * pre_local_hr_factor * ms_to_sample))
        q_samples = (round(clinical_q_range[0] * pre_local_hr_factor * ms_to_sample), 
                     round(clinical_q_range[1] * pre_local_hr_factor * ms_to_sample))
        s_samples = (round(clinical_s_range[0] * post_local_hr_factor * ms_to_sample), 
                     round(clinical_s_range[1] * post_local_hr_factor * ms_to_sample))
        t_samples = (round(clinical_t_range[0] * post_local_hr_factor * ms_to_sample), 
                     round(clinical_t_range[1] * post_local_hr_factor * ms_to_sample))
        
        # 피크 검출
        q_idx = find_peak_in_range(signal, r_idx, q_samples, min_mode=True, tol=tol)
        if q_idx is not None and q_idx < next_r_idx: 
            peaks['ECG_Q_Peaks'].append(q_idx)
        
        s_idx = find_peak_in_range(signal, r_idx, s_samples, min_mode=True, tol=tol)
        if s_idx is not None and s_idx < next_r_idx: 
            peaks['ECG_S_Peaks'].append(s_idx)
        
        t_idx = find_peak_in_range(signal, r_idx, t_samples, min_mode=False, tol=tol)
        if t_idx is not None and t_idx < next_q_idx: 
            peaks['ECG_T_Peaks'].append(t_idx)
        
        p_idx = find_peak_in_range(signal, r_idx, p_samples, min_mode=False, tol=tol)
        if p_idx is not None and p_idx < next_r_idx: 
            peaks['ECG_P_Peaks'].append(p_idx)
    
        # 디버그 정보 출력 (필요한 경우만)
        if debug and (pre_local_hr > 100 or post_local_hr > 100 or post_local_hr < 60 or pre_local_hr < 60):
            print(f"R-peak 위치: {r_idx}")
            print(f"이전 심박수: {pre_local_hr:.1f} bpm, 이전 RR 간격: {pre_local_rr:.1f} ms")
            print(f"다음 심박수: {post_local_hr:.1f} bpm, 다음 RR 간격: {post_local_rr:.1f} ms")
            print(f"이전 심박수 범위: {pre_local_hr_factor:.2f}, 다음 심박수 범위: {post_local_hr_factor:.2f}")
            print("범위 정보 (sample):")
            print(f"P-wave 범위: {r_idx + p_samples[0]} ~ {r_idx + p_samples[1]} (샘플: {p_samples[0]} ~ {p_samples[1]})")
            print(f"Q-wave 범위: {r_idx + q_samples[0]} ~ {r_idx + q_samples[1]} (샘플: {q_samples[0]} ~ {q_samples[1]})")
            print(f"S-wave 범위: {r_idx + s_samples[0]} ~ {r_idx + s_samples[1]} (샘플: {s_samples[0]} ~ {s_samples[1]})")
            print(f"T-wave 범위: {r_idx + t_samples[0]} ~ {r_idx + t_samples[1]} (샘플: {t_samples[0]} ~ {t_samples[1]})")
            print(f"P-peak 인덱스: {p_idx}")
            print(f"Q-peak 인덱스: {q_idx}")
            print(f"R-peak 인덱스: {r_idx}")
            print(f"S-peak 인덱스: {s_idx}")
            print(f"T-peak 인덱스: {t_idx}")
            print("-" * 50)
            
    # 배열로 변환 (한 번에 처리)
    for key in peaks:
        if key != 'ECG_R_Peaks' and peaks[key]:  # 빈 리스트가 아닌 경우만 변환
            peaks[key] = np.array(peaks[key])
    
    return peaks


def find_peak_in_range(signal, r_idx, sample_range, min_mode=False, tol=0.2):
    """
    지정된 범위 내에서 최대값 또는 최소값을 찾는 함수
    
    Args:
        signal: ECG 신호
        r_idx: 기준 R-peak 위치
        sample_range: 탐색 범위 (시작, 끝) 튜플
        min_mode: True면 최소값, False면 최대값 탐색
        
    Returns:
        peak_idx: 찾은 피크의 인덱스, 유효한 범위가 없으면 None
    """
    start = max(0, round(r_idx + sample_range[0] * tol + sample_range[0]))
    end = min(len(signal) - 1, round(r_idx + sample_range[1] * tol + sample_range[1]))
    
    if start >= end:
        return None
    
    if min_mode:
        peak_offset = np.argmin(signal[start:end])
    else:
        peak_offset = np.argmax(signal[start:end])
    
    return start + peak_offset

# 샘플을 밀리초(ms)로 변환하는 함수 정의
def samples_to_ms(samples, sampling_rate):
    """
    샘플 단위를 밀리초(ms) 단위로 변환하는 함수
    
    Args:
        samples: 샘플 단위의 값 또는 값들의 리스트/배열
        sampling_rate: 샘플링 레이트 (Hz)
        
    Returns:
        ms: 밀리초 단위로 변환된 값 또는 값들의 리스트/배열
    """
    return np.array(samples) * (1000 / sampling_rate)

def calculate_statistics(data, feature_name):
    """
    데이터의 통계 정보를 계산하고 출력하는 함수
    
    Args:
        data: 통계를 계산할 데이터 리스트
        feature_name: 특성 이름 (출력용)
        
    Returns:
        stats: 통계 정보를 담은 딕셔너리
    """
    # 통계 정보 계산
    avg_value = np.mean(data)
    std_value = np.std(data)
    min_value = np.min(data)
    max_value = np.max(data)
    median_value = np.median(data)
    q1_value = np.percentile(data, 25)
    q3_value = np.percentile(data, 75)
    iqr_value = q3_value - q1_value

    # 통계 정보를 딕셔너리로 반환
    stats = {
        "mean": avg_value,
        "std": std_value,
        "min": min_value,
        "max": max_value,
        "median": median_value,
        "q1": q1_value,
        "q3": q3_value,
        "iqr": iqr_value
    }
    
    return stats

def calculate_tinn(rr_intervals, bin_width=7.8125):
    """
    TINN 계산 함수
    
    Args:
        rr_intervals: RR 간격 배열 (밀리초 단위)
        bin_width: 히스토그램 bin의 너비 (ms)
        
    Returns:
        tinn: TINN 값 (ms)
        triangle_points: (M, N, O) 삼각형의 세 점
    """
    if np.min(rr_intervals) == np.max(rr_intervals):
        print("경고: 모든 RR 간격이 동일합니다. TINN을 0으로 설정합니다.")
        return 0, (np.min(rr_intervals), np.min(rr_intervals), np.min(rr_intervals))
    # 히스토그램 생성
    bins = np.arange(min(rr_intervals), max(rr_intervals) + bin_width, bin_width)
    hist, bin_edges = np.histogram(rr_intervals, bins=bins)
    # 최대 빈도 지점 찾기 (N)
    max_idx = np.argmax(hist)
    n_x = (bin_edges[max_idx] + bin_edges[max_idx + 1]) / 2
    n_y = hist[max_idx]
    
    # 삼각형 맞추기 (최적의 M과 O 찾기)
    best_fit = float('inf')
    m, o = 0, 0
    
    for i in range(max_idx):
        m_candidate = bin_edges[i]
        for j in range(max_idx + 1, len(bin_edges) - 1):
            o_candidate = bin_edges[j]
            
            # 삼각형 모델 생성
            triangle = np.zeros_like(hist)
            for k in range(len(hist)):
                bin_center = (bin_edges[k] + bin_edges[k + 1]) / 2
                if bin_center <= n_x:
                    # 왼쪽 부분
                    if bin_center >= m_candidate:
                        triangle[k] = n_y * (bin_center - m_candidate) / (n_x - m_candidate)
                else:
                    # 오른쪽 부분
                    if bin_center <= o_candidate:
                        triangle[k] = n_y * (o_candidate - bin_center) / (o_candidate - n_x)
            
            # 오차 계산
            error = np.sum((hist - triangle) ** 2)
            if error < best_fit:
                best_fit = error
                m, o = m_candidate, o_candidate
    
    tinn = o - m
    triangle_points = (m, n_x, o)
    
    return tinn, triangle_points

# ECG 특성 추출 통합 함수 정의
def extract_ecg_features(peaks, sampling_rate=500):
    """
    ECG 신호에서 다양한 특성(PR 간격, QRS 지속 시간, ST 세그먼트, QT 간격)을 추출하는 통합 함수
    
    Args:
        peaks: ECG 신호에서 검출된 피크 정보를 담은 딕셔너리
        sampling_rate: 샘플링 레이트 (Hz), 기본값 500Hz
        
    Returns:
        features: 추출된 모든 특성과 통계 정보를 담은 딕셔너리
    """
    features = {}
    
    # 필요한 피크 데이터 추출
    p_peaks = peaks.get('ECG_P_Peaks', [])
    q_peaks = peaks.get('ECG_Q_Peaks', [])
    r_peaks = peaks.get('ECG_R_Peaks', [])
    s_peaks = peaks.get('ECG_S_Peaks', [])
    t_peaks = peaks.get('ECG_T_Peaks', [])
    
    rr_intervals = np.diff(r_peaks) * (1000 / sampling_rate)
    mean_rr = rr_intervals  # 평균 RR 간격 (ms)
    
    # 1. PR 간격 계산
    pr_intervals_ms = []
    min_length = min(len(p_peaks), len(r_peaks))
    for i in range(min_length):
        pr_interval = r_peaks[i] - p_peaks[i]
        pr_interval = samples_to_ms(pr_interval, sampling_rate)
        pr_intervals_ms.append(pr_interval)
    
    # 2. QRS 지속 시간 계산
    qrs_durations_ms = []
    min_length = min(len(q_peaks), len(s_peaks))
    for i in range(min_length):
        qrs_duration = s_peaks[i] - q_peaks[i]
        qrs_duration = samples_to_ms(qrs_duration, sampling_rate)
        qrs_durations_ms.append(qrs_duration)
    
    # 3. ST 세그먼트 계산
    st_segments_ms = []
    min_length = min(len(s_peaks), len(t_peaks))
    for i in range(min_length):
        st_segment = t_peaks[i] - s_peaks[i]
        st_segment = samples_to_ms(st_segment, sampling_rate)
        st_segments_ms.append(st_segment)
    
    # 4. QT 간격 계산
    qt_intervals_ms = []
    min_length = min(len(q_peaks), len(t_peaks))
    for i in range(min_length):
        qt_interval = t_peaks[i] - q_peaks[i]
        qt_interval = samples_to_ms(qt_interval, sampling_rate)
        qt_intervals_ms.append(qt_interval)
        
    # 5. QTC 간격 계산
    qtc_intervals_ms = []
    min_length = len(qt_intervals_ms)
    mean_rr_value = np.mean(mean_rr)
    for i in range(min_length):
        # 마지막 간격에는 mean_rr의 평균값 사용
        qtc_interval = qt_intervals_ms[i] / np.sqrt(mean_rr_value)
        qtc_intervals_ms.append(qtc_interval)
    
    # 각 특성에 대한 통계 계산 및 저장
    features['PR_Interval'] = {
        'values': pr_intervals_ms,
        'stats': calculate_statistics(pr_intervals_ms, "PR 간격")
    }
    
    features['QRS_Duration'] = {
        'values': qrs_durations_ms,
        'stats': calculate_statistics(qrs_durations_ms, "QRS 지속 시간")
    }
    
    features['ST_Segments'] = {
        'values': st_segments_ms,
        'stats': calculate_statistics(st_segments_ms, "ST 세그먼트")
    }
    
    features['QT_Interval'] = {
        'values': qt_intervals_ms,
        'stats': calculate_statistics(qt_intervals_ms, "QT 간격")
    }
    
    features['QTc'] = {
        'values': qtc_intervals_ms,
        'stats': calculate_statistics(qtc_intervals_ms, "QTc 간격")
    }
    return features


def calculate_hrv_features(r_peaks, sampling_rate=500):
    """
    R-peak 위치를 기반으로 HRV 특성을 계산하는 함수
    
    Args:
        r_peaks: R-peak의 위치 배열 (샘플 단위)
        sampling_rate: 샘플링 레이트 (Hz)
        
    Returns:
        hrv_features: HRV 특성을 담은 딕셔너리
    """
    # RR 간격 계산 (밀리초 단위)
    rr_intervals = np.diff(r_peaks) * (1000 / sampling_rate)
    # 시간 영역 HRV 특성
    time_domain = {}
    
    # 기본 통계
    time_domain['mean_rr'] = np.mean(rr_intervals)  # 평균 RR 간격 (ms)
    time_domain['sdnn'] = np.std(rr_intervals)      # RR 간격의 표준편차 (ms)
    
    # Heart rate 계산 (60000 / RR 간격(ms))
    heart_rates = 60000 / rr_intervals
    time_domain['hr_max_min'] = np.max(heart_rates) - np.min(heart_rates)  # 최대 심박수 - 최소 심박수
    
    # SDRR: RR 간격의 표준편차 (SDNN과 동일하지만 다른 명칭으로도 사용됨)
    time_domain['sdrr'] = np.std(rr_intervals)
    

    # NN50: 연속된 RR 간격의 차이가 50ms를 초과하는 쌍의 수
    nn50 = sum(abs(np.diff(rr_intervals)) > 50)
    time_domain['nn50'] = nn50
    
    # pNN50: NN50을 총 RR 간격 수로 나눈 비율 (%)
    if len(rr_intervals) > 1:
        time_domain['pnn50'] = (nn50 / (len(rr_intervals) - 1)) * 100
    else:
        time_domain['pnn50'] = 0
        
    # NN20: 연속된 RR 간격의 차이가 20ms를 초과하는 쌍의 수
    nn20 = sum(abs(np.diff(rr_intervals)) > 20)
    time_domain['nn20'] = nn20
    
    # pNN20: NN20을 총 RR 간격 수로 나눈 비율 (%)
    if len(rr_intervals) > 1:
        time_domain['pnn20'] = (nn20 / (len(rr_intervals) - 1)) * 100
    else:
        time_domain['pnn20'] = 0
    
    # RMSSD: 연속된 RR 간격 차이의 제곱 평균의 제곱근
    if len(rr_intervals) > 1:
        time_domain['rmssd'] = np.sqrt(np.mean(np.diff(rr_intervals) ** 2))
    else:
        time_domain['rmssd'] = 0
    
    # TINN: NN 간격(정상 RR 간격) 히스토그램의 기하학적 특성을 측정
    # https://ekja.org/upload/media/kja-22324-supplementary-material.pdf
    tinn = calculate_tinn(rr_intervals)
    time_domain['tinn'] = tinn
    
    # 비선형 HRV 특성 (Poincaré plot)
    nonlinear = {}
    
    if len(rr_intervals) > 1:
        # SD1, SD2 계산 (Poincaré plot)
        rr_n = rr_intervals[:-1]
        rr_n1 = rr_intervals[1:]
        
        sd1 = np.std(np.subtract(rr_n, rr_n1) / np.sqrt(2))
        sd2 = np.std(np.add(rr_n, rr_n1) / np.sqrt(2))
        
        nonlinear['sd1'] = sd1
        nonlinear['sd2'] = sd2
        nonlinear['sd1_sd2_ratio'] = sd2 / sd1 if sd1 > 0 else 0
    else:
        nonlinear['sd1'] = None
        nonlinear['sd2'] = None
        nonlinear['sd1_sd2_ratio'] = None
    
    # 모든 HRV 특성을 하나의 딕셔너리로 통합
    hrv_features = {
        'rr_intervals': rr_intervals.tolist(),
        'time_domain': time_domain,
        'nonlinear': nonlinear
    }
    
    return hrv_features

def visualize_hrv_features(ecg_signal, peaks, hrv_features, ecg_features=None, title=None, path=None):
    """
    ECG 신호와 HRV 특성을 종합적으로 시각화하는 함수
    
    Args:
        ecg_signal: ECG 신호 데이터
        peaks: 검출된 피크 정보를 담은 딕셔너리
        hrv_features: HRV 특성 정보를 담은 딕셔너리
        ecg_features: ECG 특성 정보를 담은 딕셔너리 (선택적)
        sampling_rate: 샘플링 레이트 (Hz)
        title: 그래프 제목 (선택적)
        path: 그래프를 저장할 경로 (선택적)
    """
    # 서브플롯 설정 (ECG 특성 유무에 따라 행 수 조정)
    n_rows = 4 if ecg_features else 3
    fig = plt.figure(figsize=(20, n_rows * 5))
    
    # 1. ECG 신호와 피크 표시
    ax1 = plt.subplot(n_rows, 1, 1)
    ax1.plot(ecg_signal, label="ECG Signal")
    
    # 피크 타입 정의
    peak_types = {
        'ECG_R_Peaks': {'color': 'red', 'marker': 'o', 'label': 'R-peaks'},
        'ECG_P_Peaks': {'color': 'orange', 'marker': 'o', 'label': 'P-peaks'},
        'ECG_Q_Peaks': {'color': 'green', 'marker': 'o', 'label': 'Q-peaks'},
        'ECG_S_Peaks': {'color': 'blue', 'marker': 'o', 'label': 'S-peaks'},
        'ECG_T_Peaks': {'color': 'purple', 'marker': 'o', 'label': 'T-peaks'}
    }
    
    # 모든 피크 타입 플롯
    for peak_name, peak_info in peak_types.items():
        if peak_name in peaks.keys() and len(peaks[peak_name]) > 0:
            ax1.scatter(
                peaks[peak_name], 
                ecg_signal[peaks[peak_name]],
                color=peak_info['color'], 
                marker=peak_info['marker'], 
                label=peak_info['label']
            )
    
    ax1.set_title("ECG Signal with Detected Peaks")
    ax1.set_ylabel("Amplitude")
    ax1.grid(True, alpha=0.3)
    ax1.legend(loc='upper left')
    
    # 2. RR 간격 시계열 플롯
    ax2 = plt.subplot(n_rows, 1, 2)
    rr_intervals = hrv_features['rr_intervals']
    ax2.plot(rr_intervals, 'o-', color='blue')
    ax2.axhline(y=np.mean(rr_intervals), color='r', linestyle='--', label=f'Mean: {np.mean(rr_intervals):.2f} ms')
    ax2.fill_between(range(len(rr_intervals)), 
                    np.mean(rr_intervals) - np.std(rr_intervals), 
                    np.mean(rr_intervals) + np.std(rr_intervals), 
                    color='red', alpha=0.2, label=f'SD: {np.std(rr_intervals):.2f} ms')
    ax2.set_title("RR Intervals Time Series")
    ax2.set_xlabel("Beat Number")
    ax2.set_ylabel("RR Interval (ms)")
    ax2.grid(True, alpha=0.3)
    ax2.legend(loc='upper left')
    
    # 3. Poincaré 플롯 (비선형 HRV 분석)
    ax3 = plt.subplot(n_rows, 1, 3)
    if len(rr_intervals) > 1:
        rr_n = np.array(rr_intervals[:-1])
        rr_n1 = np.array(rr_intervals[1:])
        
        # 산점도 그리기
        ax3.scatter(rr_n, rr_n1, color='blue', alpha=0.7)
        
        # SD1, SD2 타원 그리기
        sd1 = hrv_features['nonlinear']['sd1']
        sd2 = hrv_features['nonlinear']['sd2']
        
        if sd1 is not None and sd2 is not None:
            # 타원 중심점 (평균 RR, 평균 RR)
            center = (np.mean(rr_n), np.mean(rr_n1))
            
            # 타원 그리기
            from matplotlib.patches import Ellipse
            ellipse = Ellipse(xy=center, width=2*sd2, height=2*sd1, 
                             angle=45, edgecolor='red', fc='none', lw=2)
            ax3.add_patch(ellipse)
            
            # SD1, SD2 선 그리기
            ax3.plot([center[0], center[0]], [center[1], center[1] + sd1], 'r-', label=f'SD1: {sd1:.2f} ms')
            ax3.plot([center[0], center[0] + sd2], [center[1], center[1]], 'g-', label=f'SD2: {sd2:.2f} ms')
            
            # 정체선 그리기 (y=x)
            min_rr = min(min(rr_n), min(rr_n1))
            max_rr = max(max(rr_n), max(rr_n1))
            ax3.plot([min_rr, max_rr], [min_rr, max_rr], 'k--', alpha=0.5)
            
            # 비율 표시
            ratio = hrv_features['nonlinear']['sd1_sd2_ratio']
            ax3.text(0.05, 0.05, f'SD1/SD2: {ratio:.3f}', transform=ax3.transAxes, 
                    verticalalignment='bottom', bbox=dict(boxstyle='round', facecolor='white', alpha=0.5))
    
    ax3.set_title("Poincaré Plot")
    ax3.set_xlabel("RR(n) (ms)")
    ax3.set_ylabel("RR(n+1) (ms)")
    ax3.grid(True, alpha=0.3)
    ax3.legend(loc='upper left')
    
    # 4. ECG 특성 요약 (선택적)
    if ecg_features:
        ax4 = plt.subplot(n_rows, 1, 4)
        
        # 특성 이름과 평균값 추출
        feature_names = []
        feature_values = []
        
        for feature_name, feature_data in ecg_features.items():
            if 'stats' in feature_data and 'mean' in feature_data['stats']:
                feature_names.append(feature_name.replace('_', ' ').title())
                feature_values.append(feature_data['stats']['mean'])
        
        # 바 차트로 표시
        bars = ax4.bar(feature_names, feature_values, color=['blue', 'green', 'orange', 'purple', 'red'])
        
        # 값 표시
        for bar in bars:
            height = bar.get_height()
            ax4.text(bar.get_x() + bar.get_width()/2., height + 1,
                    f'{height:.1f} ms', ha='center', va='bottom')
        
        ax4.set_title("ECG Features Summary (Mean Values)")
        ax4.set_ylabel("Duration (ms)")
        ax4.grid(True, alpha=0.3, axis='y')
        
    # HRV 시간 영역 특성 표시
    time_domain = hrv_features['time_domain']
    info_text = (
        f"Mean HR: {60000 / np.mean(hrv_features['rr_intervals']):.2f} bpm\n"
        f"Time Domain HRV Metrics:\n"
        f"Mean RR: {time_domain['mean_rr']:.2f} ms\n"
        f"SDNN: {time_domain['sdnn']:.2f} ms\n"
        f"RMSSD: {time_domain['rmssd']:.2f} ms\n"
        f"pNN20: {time_domain['pnn20']:.2f}%\n"
        f"pNN50: {time_domain['pnn50']:.2f}%\n"
        f"HR Max-Min: {time_domain['hr_max_min']:.2f} bpm"
    )
    
    # 정보 텍스트 추가
    fig.text(0.02, 0.02, info_text, fontsize=10, 
             bbox=dict(boxstyle='round', facecolor='white', alpha=0.8))
    
    # 전체 제목 설정
    if title:
        plt.suptitle(title, fontsize=16)
    
    plt.tight_layout()
    plt.subplots_adjust(top=0.92 if title else 0.95)  # 제목 공간 확보
    
    # 저장 또는 표시
    if path:
        plt.savefig(path)
    else:
        plt.show()
    plt.close()
    
def save_images_to_pptx(image_folder, save_path, output_pptx):
    """
    이미지 폴더의 모든 이미지를 PowerPoint 문서로 저장합니다.
    각 슬라이드에 하나의 이미지만 표시됩니다.
    
    Args:
        image_folder (str): 이미지가 저장된 폴더 경로
        output_pptx (str): 저장할 PowerPoint 파일 경로
    """
    # 새 프레젠테이션 생성
    prs = Presentation()
    
    # 이미지 파일 목록 가져오기 (png 파일만)
    image_files = glob.glob(os.path.join(image_folder, "*.png"))
    
    # 파일명에서 인덱스 추출하여 정렬
    def get_index(file_path):
        # 파일명에서 마지막 언더스코어 이후의 숫자를 인덱스로 사용
        filename = os.path.basename(file_path)
        index_str = filename.split('_')[-2].split('.')[0]
        return int(index_str)
    
    # 인덱스 기준으로 이미지 파일 정렬
    image_files.sort(key=get_index)
    
    # 각 이미지를 슬라이드에 추가
    for image_path in image_files:
        # 빈 슬라이드 추가
        slide_layout = prs.slide_layouts[6]  # 빈 레이아웃
        slide = prs.slides.add_slide(slide_layout)
        
        # 슬라이드 크기 가져오기
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # 여백 설정
        margin = Inches(0.5)
        
        # 이미지를 슬라이드 크기에 맞게 조정
        left = margin
        top = margin
        width = slide_width - (2 * margin)
        height = slide_height - (2 * margin) - Inches(0.5)  # 텍스트 상자 공간 확보
        
        
        # 이미지 추가 (슬라이드 크기에 맞게 자동 조정)
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    
    # PowerPoint 파일 저장
    prs.save(os.path.join(save_path, output_pptx))
    print(f"PowerPoint 문서가 {output_pptx}에 저장되었습니다.")
    
def assess_signal_quality(signal, sampling_rate):
    """
    ECG 신호의 품질을 평가하는 함수
    
    Parameters:
    -----------
    signal : numpy.ndarray
        ECG 신호 데이터
    sampling_rate : float
        샘플링 레이트
        
    Returns:
    --------
    quality_score : float
        신호 품질 점수 (0~1 사이 값, 1이 최상)
    """
    # 신호가 없거나 모두 0인 경우
    if len(signal) == 0 or np.all(signal == 0):
        return 0.0
    
    # 1. 신호 대 잡음비(SNR) 추정
    # 신호의 표준편차와 중앙값 절대편차(MAD)를 사용하여 SNR 추정
    signal_std = np.std(signal)
    signal_mad = np.median(np.abs(signal - np.median(signal)))
    estimated_snr = signal_std / (signal_mad + 1e-10)  # 0 나누기 방지
    snr_score = min(1.0, estimated_snr / 5.0)  # SNR 5 이상이면 최대 점수
    
    # 2. 기저선 변동 평가
    # 저주파 성분 추출을 위한 필터링
    nyquist = 0.5 * sampling_rate
    low_freq = 0.5 / nyquist
    b, a = scipy.signal.butter(2, low_freq, 'lowpass')
    baseline = scipy.signal.filtfilt(b, a, signal)
    baseline_variation = np.std(baseline) / (np.std(signal) + 1e-10)
    baseline_score = 1.0 - min(1.0, baseline_variation)
    
    # 3. 이상치 비율 평가
    # 평균에서 3 표준편차 이상 벗어난 샘플 비율
    mean = np.mean(signal)
    std = np.std(signal)
    outliers = np.abs(signal - mean) > 3 * std
    outlier_ratio = np.sum(outliers) / len(signal)
    outlier_score = 1.0 - min(1.0, outlier_ratio * 10)  # 10% 이상이면 최저 점수
    
    # 4. 신호 연속성 평가
    # 연속된 샘플 간의 급격한 변화 감지
    diff = np.abs(np.diff(signal))
    sudden_changes = diff > 5 * np.std(diff)
    continuity_score = 1.0 - min(1.0, np.sum(sudden_changes) / len(diff) * 10)
    
    # 종합 품질 점수 계산 (가중 평균)
    quality_score = (snr_score * 0.4 + 
                     baseline_score * 0.3 + 
                     outlier_score * 0.2 + 
                     continuity_score * 0.1)
    
    return quality_score

def extract_qrs_complexes(ecg_signal, qrs_mask, fixed_length=32, r_peaks=None):
    """
    ECG 신호에서 QRS 컴플렉스를 추출하고 고정 길이로 리샘플링
    주어진 R-peak 위치도 함께 리샘플링
    
    Args:
        ecg_signal: ECG 신호
        qrs_mask: QRS 컴플렉스 마스크
        fixed_length: 리샘플링할 고정 길이
        r_peaks: 원본 신호에서의 R-peak 위치 리스트 (선택적)
    """
    # 텐서를 numpy 배열로 변환
    if isinstance(ecg_signal, torch.Tensor):
        ecg_signal = ecg_signal.cpu().numpy()
    if isinstance(qrs_mask, torch.Tensor):
        qrs_mask = qrs_mask.cpu().numpy()
    
    # 차원 확인 및 조정
    if len(ecg_signal.shape) > 1:
        ecg_signal = ecg_signal.squeeze()
    if len(qrs_mask.shape) > 1:
        qrs_mask = qrs_mask.squeeze()
    
    # QRS 컴플렉스 구간 찾기
    qrs_regions = []
    in_qrs = False
    start_idx = 0
    
    for i, val in enumerate(qrs_mask):
        if val > 0.5 and not in_qrs:  # QRS 시작
            in_qrs = True
            start_idx = i
        elif val <= 0.5 and in_qrs:  # QRS 종료
            in_qrs = False
            qrs_regions.append((start_idx, i))
    # 마지막 QRS가 신호 끝까지 계속되는 경우
    if in_qrs:
        qrs_regions.append((start_idx, len(qrs_mask)))
    # QRS 컴플렉스 추출 및 리샘플링
    resampled_qrs = []
    resampled_r_peaks = []
    
    for i, (start, end) in enumerate(qrs_regions):
        if end - start < 5:  # 너무 짧은 QRS는 무시
            continue
        qrs_complex = ecg_signal[start:end]
        
        # 리샘플링 (고정 길이로)
        resampled = scipy.signal.resample(qrs_complex, fixed_length)
        resampled_qrs.append(resampled)
        
        # R-peak 위치 리샘플링
        if r_peaks is not None and i < len(r_peaks):
            # 주어진 R-peak이 현재 QRS 영역 내에 있는지 확인
            r_peak = r_peaks[i]
            if start <= r_peak < end:
                # 원본 영역 내에서의 상대적 위치 계산
                relative_pos = (r_peak - start) / (end - start)
                # 리샘플링된 영역에서의 위치 계산
                resampled_r_peak = int(relative_pos * fixed_length)
                resampled_r_peaks.append(resampled_r_peak)
        else:
            # R-peak이 주어지지 않은 경우 신호의 최대값 위치로 가정
            original_r_peak = np.argmax(qrs_complex)
            resampled_r_peak = int(original_r_peak * (fixed_length / len(qrs_complex)))
            resampled_r_peaks.append(resampled_r_peak)
    
    return np.array(resampled_qrs), qrs_regions, np.array(resampled_r_peaks)

def refine_qrs_mask(qrs_mask, min_qrs_width=35, max_qrs_width=60, min_gap=60):
    """
    QRS 마스크를 후처리하여 너무 좁은 QRS 영역 제거, 너무 넓은 QRS 영역 분할, 
    그리고 너무 가까운 QRS 영역 병합
    
    Args:
        qrs_mask: 원본 QRS 마스크
        min_qrs_width: 최소 QRS 너비 (샘플 수)
        max_qrs_width: 최대 QRS 너비 (샘플 수)
        min_gap: QRS 영역 간 최소 간격 (샘플 수)
    
    Returns:
        refined_mask: 정제된 QRS 마스크
    """
    # 마스크 복사
    refined_mask = qrs_mask.copy()
    
    # QRS 영역 찾기
    qrs_regions = []
    in_qrs = False
    start_idx = 0
    
    for i, val in enumerate(refined_mask):
        if val > 0.5 and not in_qrs:  # QRS 시작
            in_qrs = True
            start_idx = i
        elif val <= 0.5 and in_qrs:  # QRS 종료
            in_qrs = False
            qrs_regions.append((start_idx, i))
    
    # 마지막 QRS가 신호 끝까지 계속되는 경우
    if in_qrs:
        qrs_regions.append((start_idx, len(refined_mask)))
    
    # 너무 좁은 QRS 영역 제거
    valid_regions = []
    for start, end in qrs_regions:
        if end - start >= min_qrs_width:
            valid_regions.append((start, end))
        else:
            refined_mask[start:end] = 0  # 마스크에서 제거
    
    # 너무 넓은 QRS 영역 분할
    split_regions = []
    for start, end in valid_regions:
        width = end - start
        if width > max_qrs_width:
            # 신호 기반으로 분할점 찾기 (여기서는 간단히 중간점으로 분할)
            mid = start + width // 2
            split_regions.append((start, mid))
            split_regions.append((mid, end))
        else:
            split_regions.append((start, end))
    
    # 너무 가까운 QRS 영역 병합
    merged_regions = []
    if split_regions:
        current_start, current_end = split_regions[0]
        
        for i in range(1, len(split_regions)):
            next_start, next_end = split_regions[i]
            
            if next_start - current_end < min_gap:
                # 병합
                current_end = next_end
            else:
                merged_regions.append((current_start, current_end))
                current_start, current_end = next_start, next_end
        
        # 마지막 영역 추가
        merged_regions.append((current_start, current_end))
    
    # 정제된 마스크 생성
    refined_mask = np.zeros_like(qrs_mask)
    for start, end in merged_regions:
        refined_mask[start:end] = 1
    
    return refined_mask

def validate_r_peaks(ecg_signal, r_peak_locations, fs=500, min_rr_ms=200, max_rr_ms=2000):
    """
    생리학적 제약 조건을 기반으로 R-peak 위치 검증 및 보정
    
    Args:
        ecg_signal: ECG 신호
        r_peak_locations: 검출된 R-peak 위치 목록
        fs: 샘플링 주파수 (Hz)
        min_rr_ms: 최소 R-R 간격 (밀리초)
        max_rr_ms: 최대 R-R 간격 (밀리초)
    
    Returns:
        validated_peaks: 검증된 R-peak 위치 목록
    """
    if len(r_peak_locations) <= 1:
        return r_peak_locations
    
    # 샘플 단위로 변환
    min_rr_samples = int(min_rr_ms * fs / 1000)
    max_rr_samples = int(max_rr_ms * fs / 1000)
    
    # R-peak 위치 정렬
    r_peak_locations = sorted(r_peak_locations)
    
    # 검증된 R-peak 목록 초기화
    validated_peaks = [r_peak_locations[0]]  # 첫 번째 피크는 유지
    
    for i in range(1, len(r_peak_locations)):
        current_peak = r_peak_locations[i]
        last_valid_peak = validated_peaks[-1]
        rr_interval = current_peak - last_valid_peak
        
        # 최소 R-R 간격 검사
        if rr_interval < min_rr_samples:
            # 두 피크 중 더 높은 진폭을 가진 피크 선택
            if current_peak < len(ecg_signal) and last_valid_peak < len(ecg_signal):
                if abs(ecg_signal[current_peak]) > abs(ecg_signal[last_valid_peak]):
                    # 현재 피크가 더 강함 - 이전 피크 대체
                    validated_peaks[-1] = current_peak
            continue
        
        # 최대 R-R 간격 검사 (선택적)
        if rr_interval > max_rr_samples:
            # 놓친 R-peak가 있을 수 있음 - 중간에 피크 찾기 시도
            mid_point = (last_valid_peak + current_peak) // 2
            window_size = min(rr_interval // 4, 100)  # 검색 윈도우 크기
            
            # 중간 지점 주변에서 피크 찾기
            search_start = max(0, mid_point - window_size)
            search_end = min(len(ecg_signal), mid_point + window_size)
            
            if search_end > search_start:
                # 검색 윈도우 내에서 가장 큰 피크 찾기
                window_signal = np.abs(ecg_signal[search_start:search_end])
                max_idx = np.argmax(window_signal)
                potential_peak = search_start + max_idx
                
                # 충분히 큰 피크인지 확인
                if window_signal[max_idx] > 0.5 * abs(ecg_signal[last_valid_peak]):
                    validated_peaks.append(potential_peak)
        
        # 정상 R-R 간격 - 피크 추가
        validated_peaks.append(current_peak)
    
    return validated_peaks

def refine_r_peak_locations(ecg_signal, r_peak_locations, window_size=10):
    """
    R-peak 위치를 주변 윈도우에서 실제 최대값으로 정밀화
    
    Args:
        ecg_signal: ECG 신호
        r_peak_locations: 초기 R-peak 위치 목록
        window_size: 검색 윈도우 크기 (샘플 수)
    
    Returns:
        refined_peaks: 정밀화된 R-peak 위치 목록
    """
    refined_peaks = []
    
    for peak in r_peak_locations:
        # 검색 윈도우 범위 설정
        window_start = max(0, peak - window_size)
        window_end = min(len(ecg_signal), peak + window_size + 1)
        
        # 윈도우 내에서 최대값 위치 찾기
        window = ecg_signal[window_start:window_end]
        
        # 절대값 기준으로 최대값 찾기 (양수/음수 피크 모두 고려)
        local_max_idx = np.argmax(window)
        
        # 전체 신호에서의 위치로 변환
        refined_peak = window_start + local_max_idx
        refined_peaks.append(refined_peak)
    
    return refined_peaks

def r_peak_detection(ecg_signal, qrs_mask, rpeak_model, device, fs=500):
    """
    R-peak 검출 함수
    
    Args:
        ecg_signal: ECG 신호
        qrs_mask: QRS 마스크
        rpeak_model: R-peak 위치 예측 회귀 모델
        device: 연산 장치 (CPU 또는 GPU)
        fs: 샘플링 주파수 (Hz)
    
    Returns:
        r_peak_locations: 검출된 R-peak 위치 목록
    """
    # 1. QRS 마스크 후처리
    refined_mask = refine_qrs_mask(qrs_mask)
    # 2. QRS 컴플렉스 추출
    qrs_complexes, qrs_regions, _ = extract_qrs_complexes(ecg_signal, refined_mask)
    if len(qrs_complexes) == 0:
        return []
    # 3. 각 QRS 컴플렉스에서 R 피크 위치 예측
    r_peak_locations = []
    with torch.no_grad():
        for i, qrs in enumerate(qrs_complexes):
            # 입력 텐서 형태 확인 및 조정
            qrs_tensor = torch.FloatTensor(qrs).unsqueeze(0).unsqueeze(0).to(device)
            
            # 정규화된 위치 예측 (0~1 사이 값)
            normalized_position = rpeak_model(qrs_tensor).item()
            normalized_position = max(0, min(1, normalized_position))
            # 정규화된 위치를 실제 샘플 위치로 변환
            start, end = qrs_regions[i]
            r_peak_idx = start + int(normalized_position * (end - start))
            r_peak_locations.append(r_peak_idx)
    # 4. R-peak 위치 정밀화
    refined_peaks = refine_r_peak_locations(ecg_signal, r_peak_locations)
    
    # 5. 생리학적 제약 기반 검증
    validated_peaks = validate_r_peaks(ecg_signal, refined_peaks, fs)
    
    return validated_peaks

def predict_r_peaks(ecg_signal, qrs_mask, rpeak_model, device, fs=500):
    """
    R-peak 검출 함수
    
    Args:
        ecg_signal: ECG 신호 데이터
        qrs_mask: QRS 복합체 마스크 (U-Net 예측 결과)
        rpeak_model: R 피크 위치 예측 회귀 모델
        device: 연산 장치 (CPU 또는 GPU)
        fs: 샘플링 주파수 (Hz)
        
    Returns:
        dict: R-peak 검출 결과와 신호 품질 정보를 포함한 딕셔너리
    """
    rpeak_model.to(device)
    rpeak_model.eval()
    
    # 텐서를 numpy 배열로 변환
    if isinstance(ecg_signal, torch.Tensor):
        ecg_signal = ecg_signal.cpu().numpy()
    if isinstance(qrs_mask, torch.Tensor):
        qrs_mask = qrs_mask.cpu().numpy()
    
    # 차원 확인 및 조정
    if len(ecg_signal.shape) > 1:
        ecg_signal = ecg_signal.squeeze()
    if len(qrs_mask.shape) > 1:
        qrs_mask = qrs_mask.squeeze()
    
    # R-peak 검출
    r_peak_locations = r_peak_detection(ecg_signal, qrs_mask, rpeak_model, device, fs)
    
    # 결과 딕셔너리 생성
    result = {
        'R-peak Regression': {
            'original_rpeaks': np.array(r_peak_locations),
            'corrected_rpeaks': np.array(r_peak_locations),
            'plot_signal': ecg_signal
        }
    }
    
    return result

# 정답 값과 모델 예측 결과 비교
def r_peak_prediction(ecg_signal, unet_model, rpeak_model, device, best_threshold, sampling_rate=500):
    """
    ECG 신호에서 R-peak를 예측하는 함수
    
    Parameters:
    -----------
    ecg_signal : numpy.ndarray
        ECG 신호 데이터 (numpy 배열)
    unet_model : torch.nn.Module
        학습된 U-Net 모델
    rpeak_model : torch.nn.Module
        학습된 R-peak 회귀 모델
    device : torch.device
        연산을 수행할 디바이스
    best_threshold : float
        U-Net 모델의 출력을 이진화하기 위한 임계값
    sampling_rate : int, optional
        ECG 신호의 샘플링 주파수 (Hz), 기본값: 500
        
    Returns:
    --------
    numpy.ndarray
        최종 검출된 R-peak 위치 배열
    """
    # numpy 배열을 텐서로 변환
    if isinstance(ecg_signal, np.ndarray):
        test_ecg = torch.from_numpy(ecg_signal).float().to(device)
    else:
        test_ecg = ecg_signal.to(device)
    
    # 차원 확인 및 조정
    if len(test_ecg.shape) == 1:
        test_ecg = test_ecg.unsqueeze(0).unsqueeze(0)  # [1, 1, signal_length]
    elif len(test_ecg.shape) == 2:
        test_ecg = test_ecg.unsqueeze(1)  # [batch, 1, signal_length]
    
    # U-Net 모델로 QRS 마스크 예측
    with torch.no_grad():
        outputs = unet_model(test_ecg)
        qrs_mask = (torch.sigmoid(outputs) > best_threshold).float().cpu().squeeze().numpy()
    
    # ECG 신호를 numpy로 변환
    current_ecg_np = test_ecg.cpu().numpy().squeeze()
    
    # 앙상블 방법으로 R-peak 검출
    result = extract_ecg_and_rpeaks(current_ecg_np, sampling_rate, 0.5, 40, min_votes_ratio=0.5)
    
    # 회귀 모델로 R-peak 예측
    model_rpeaks = predict_r_peaks(current_ecg_np, qrs_mask, rpeak_model, device, sampling_rate)
    
    # 앙상블 방법으로 검출된 R-peak
    ensemble_rpeaks = np.array(result['ensemble_rpeaks'])
    
    # 회귀 모델로 예측된 R-peak
    regression_rpeaks = np.array(model_rpeaks['R-peak Regression']['corrected_rpeaks'])
    
    # 신호 품질에 따른 최종 R-peak 결정
    signal_quality = result['signal_quality']
    
    # 샘플링 레이트 기반 시간 간격 계산 (ms 단위)
    time_150ms = int(0.15 * sampling_rate)
    time_100ms = int(0.1 * sampling_rate)
    time_50ms = int(0.05 * sampling_rate)
    
    # 신호 품질에 따른 R-peak 결정 전략
    if signal_quality > 0.85:  # 높은 신호 품질
        # 앙상블 결과 우선 사용, 회귀 모델로 보완
        final_rpeaks = ensemble_rpeaks.copy()
        
        # 앙상블에서 놓친 R-peak를 회귀 모델에서 찾아 추가
        if len(ensemble_rpeaks) > 0:
            for reg_peak in regression_rpeaks:
                # 가장 가까운 앙상블 피크와의 거리 계산
                min_distance = np.min(np.abs(ensemble_rpeaks - reg_peak))
                if min_distance > time_150ms:  # 150ms 이상 차이나면 추가
                    final_rpeaks = np.append(final_rpeaks, reg_peak)
        else:
            final_rpeaks = regression_rpeaks.copy()
                
    elif signal_quality > 0.6:  # 중간 신호 품질
        # 회귀 모델 결과 우선 사용
        final_rpeaks = regression_rpeaks.copy()
        
        # 회귀 모델에 없는 앙상블 피크 추가 검토
        if len(regression_rpeaks) > 0:
            for ens_peak in ensemble_rpeaks:
                min_distance = np.min(np.abs(regression_rpeaks - ens_peak))
                if min_distance > time_100ms:  # 100ms 이상 차이나면 추가
                    final_rpeaks = np.append(final_rpeaks, ens_peak)
        else:
            final_rpeaks = ensemble_rpeaks.copy()
                
    else:  # 낮은 신호 품질
        # 두 모델 모두에서 검출된 피크만 신뢰
        final_rpeaks = np.array([], dtype=int)
        
        # 두 결과가 유사한 피크만 선택
        for peak in regression_rpeaks:
            if len(ensemble_rpeaks) > 0:
                min_distance = np.min(np.abs(ensemble_rpeaks - peak))
                if min_distance < time_50ms:  # 50ms 이내에 앙상블 피크가 있으면 추가
                    final_rpeaks = np.append(final_rpeaks, peak)
        
        # 결과가 너무 적으면 회귀 모델 결과 사용
        if len(final_rpeaks) < 3 and len(regression_rpeaks) >= 3:
            final_rpeaks = regression_rpeaks.copy()
    
    # 중복 제거 및 정렬
    if len(final_rpeaks) > 0:
        final_rpeaks = np.unique(final_rpeaks.astype(int))
        final_rpeaks = np.sort(final_rpeaks)
    
    results = {
        'regression_rpeaks': regression_rpeaks,
        'ensemble_rpeaks': final_rpeaks,
        'plot_signal': ecg_signal,
        'signal_quality': signal_quality
    }
    return results

# 모델 평가 함수 정의
def evaluate_model_unet(model, test_loader, device):
    """
    모델을 평가하고 성능 지표를 계산하는 함수
    
    Args:
        model: 평가할 모델
        test_loader: 테스트 데이터 로더
        device: 연산을 수행할 장치 (CPU 또는 GPU)
    
    Returns:
        평가 지표 (정확도, 정밀도, 재현율, F1 점수)
    """
    model.eval()
    all_preds = []
    all_labels = []
    all_outputs = []
    
    # 테스트 데이터에 대한 예측 수행
    with torch.no_grad():
        for inputs, labels, _ in tqdm(test_loader, desc="모델 평가 중"):
            inputs = inputs.to(device)
            labels = labels.to(device)
            
            outputs = model(inputs)
            probs = torch.sigmoid(outputs)
            
            # 예측 결과와 실제 레이블 수집
            all_preds.extend(probs.cpu().numpy().flatten())
            all_outputs.extend(outputs.cpu().numpy().flatten())
            all_labels.extend(labels.cpu().numpy().flatten())
    
    # 수집된 데이터를 numpy 배열로 변환
    all_preds = np.array(all_preds)
    all_labels = np.array(all_labels)
    all_outputs = np.array(all_outputs)
    
    # 최적의 임계값 찾기
    thresholds = np.arange(0.1, 0.9, 0.05)
    best_f1 = 0
    best_threshold = 0.5
    
    for threshold in tqdm(thresholds, desc="임계값 최적화 중"):
        binary_preds = (all_preds > threshold).astype(int)
        f1 = f1_score(all_labels, binary_preds, average='binary')
        
        if f1 > best_f1:
            best_f1 = f1
            best_threshold = threshold
    
    # 최적 임계값으로 최종 예측
    final_preds = (all_preds > best_threshold).astype(int)
    
    # 성능 지표 계산
    accuracy = accuracy_score(all_labels, final_preds)
    precision = precision_score(all_labels, final_preds, average='binary')
    recall = recall_score(all_labels, final_preds, average='binary')
    f1 = f1_score(all_labels, final_preds, average='binary')
    
    print(f"최적 임계값: {best_threshold:.3f}")
    print(f"정확도: {accuracy:.4f}")
    print(f"정밀도: {precision:.4f}")
    print(f"재현율: {recall:.4f}")
    print(f"F1 점수: {f1:.4f}")
    
    return accuracy, precision, recall, f1, best_threshold